"""
Restore the 19-slide base from the backup, then add Part 3 slides.
"""
import copy
import shutil
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.shapes.shapetree import SlideShapes

BACKUP_PATH = "GitHub-Pages-Training.pptx.bak"
OUTPUT_PATH  = "GitHub-Pages-Training.pptx"


# ── slide removal ─────────────────────────────────────────────────────────────

def remove_slide(prs, index):
    """Remove a slide at the given index from the presentation."""
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[index]
    prs.part.drop_rel(prs.slides._sldIdLst[index].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'))
    del xml_slides[index]


# ── run access ────────────────────────────────────────────────────────────────

def get_all_runs(slide):
    spTree = slide._element.find('.//' + qn('p:spTree'))
    shapes = SlideShapes(spTree, slide)
    runs = []
    for shape in shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    runs.append(run)
    return runs


# ── slide cloning ─────────────────────────────────────────────────────────────

def clone_and_replace(prs, source_index, new_texts):
    src = prs.slides[source_index]
    new_slide = prs.slides.add_slide(src.slide_layout)

    src_spTree = src._element.find('.//' + qn('p:spTree'))
    dst_cSld = new_slide._element.find(qn('p:cSld'))
    dst_spTree = dst_cSld.find(qn('p:spTree'))
    new_spTree = copy.deepcopy(src_spTree)
    dst_cSld.replace(dst_spTree, new_spTree)

    runs = get_all_runs(new_slide)
    for i, run in enumerate(runs):
        run.text = new_texts[i] if i < len(new_texts) else ""
    return new_slide


# ── slide content ─────────────────────────────────────────────────────────────

# Template source indices in the 19-slide deck:
#   2 → section divider  (5 runs):  PART, title, subtitle, footer×2
#   3 → concept slide   (13 runs):  label, title, pagenum, bullet×6, callout-label, callout-main, callout-example, footer
#   5 → step+code slide (19 runs):  step-label, step-title, pagenum, intro, code×12, warn-icon, warn-text, footer

DIVIDER = 2
CONCEPT  = 3
STEP     = 5

SLIDES = [
    (DIVIDER, [
        "PART 3",
        "Fork, Contribute & PR",
        "Four steps from idea to merged pull request",
        "https://teletrex.com",
        "https://teletrex.com",
    ]),
    (CONCEPT, [
        "CONCEPTS",
        "What is Forking?",
        "20",
        "A fork is your personal copy of someone else's repository",
        "Lives under your GitHub account — you have full write access",
        "The original repo is called the upstream",
        "Changes in your fork don't affect upstream until a PR is merged",
        "Pull Requests propose your changes back to the original project",
        "Maintainers review, discuss, and merge (or decline) your PR",
        "Example repo for today",
        "https://github.com/hd-admin/hackerdojo.org",
        "Your fork: https://github.com/YOUR-USERNAME/hackerdojo.org",
        "https://teletrex.com",
    ]),
    (STEP, [
        "STEP 1 OF 4",
        "Fork the repository on GitHub",
        "21",
        "Navigate to the repo on GitHub and click Fork in the top-right corner.",
        "# Open the upstream repo in your browser",
        "https://github.com/hd-admin/hackerdojo.org",
        "",
        "# Click the Fork button (top-right corner)",
        "# GitHub creates your personal copy instantly:",
        "https://github.com/YOUR-USERNAME/hackerdojo.org",
        "",
        "",
        "",
        "",
        "",
        "",
        "💡  ",
        "Your fork is fully independent. Pushes go to YOUR copy only — nothing changes upstream until you open a PR.",
        "https://teletrex.com",
    ]),
    (STEP, [
        "STEP 2 OF 4",
        "Clone your fork to your machine",
        "22",
        "Clone YOUR fork (not the original), then add upstream so you can pull future updates.",
        "# Clone your fork",
        "git clone https://github.com/YOUR-USERNAME/hackerdojo.org",
        "",
        "# Enter the project folder",
        "cd hackerdojo.org",
        "",
        "# Track the original as 'upstream'",
        "git remote add upstream https://github.com/hd-admin/hackerdojo.org",
        "",
        "# Confirm both remotes are set",
        "git remote -v",
        "",
        "💡  ",
        "Always clone YOUR fork, not the upstream. Cloning upstream gives you no push access and skips the PR workflow.",
        "https://teletrex.com",
    ]),
    (STEP, [
        "STEP 3 OF 4",
        "Create a branch, make your changes, and commit",
        "23",
        "Work on a dedicated branch — never commit directly to main.",
        "# Create and switch to a feature branch",
        "git checkout -b fix/update-homepage-text",
        "",
        "# Edit the file you want to change",
        "# e.g. open index.html in VS Code",
        "",
        "# Stage your changes",
        "git add index.html",
        "",
        "# Commit with a clear, descriptive message",
        'git commit -m "Update homepage welcome text"',
        "",
        "💡  ",
        "Branch names like fix/…, feat/…, or docs/… signal intent at a glance and keep your main branch clean.",
        "https://teletrex.com",
    ]),
    (STEP, [
        "STEP 4 OF 4",
        "Push your branch and open a Pull Request",
        "24",
        "Push the branch to YOUR fork, then open a PR against the upstream repo.",
        "# Push the branch to your fork",
        "git push -u origin fix/update-homepage-text",
        "",
        "# Open a PR with the GitHub CLI",
        "gh pr create --repo hd-admin/hackerdojo.org \\",
        '  --title "Update homepage welcome text" \\',
        '  --body "Improves the welcome message for new visitors."',
        "",
        "# GitHub also prints a URL after the push — open it to",
        "# create the PR in the browser instead.",
        "",
        "",
        "💡  ",
        "After your PR is merged: git pull upstream main && git push origin main  keeps your fork in sync.",
        "https://teletrex.com",
    ]),
]


# ── main ──────────────────────────────────────────────────────────────────────

prs = Presentation(BACKUP_PATH)
print(f"Loaded backup: {len(prs.slides)} slides")

# Strip the 6 bad slides added in earlier runs (slides 20-25, indices 19-24)
while len(prs.slides) > 19:
    remove_slide(prs, len(prs.slides) - 1)
    print(f"  Removed slide → {len(prs.slides)} remaining")

print(f"Base deck restored: {len(prs.slides)} slides")

for tmpl_idx, texts in SLIDES:
    slide = clone_and_replace(prs, tmpl_idx, texts)
    print(f"  Added: {texts[1]!r}")

print(f"Final slide count: {len(prs.slides)}")
prs.save(OUTPUT_PATH)
print(f"Saved → {OUTPUT_PATH}")
